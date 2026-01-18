from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from datetime import datetime

class PowerPointGenerator:
    """PowerPointレポート生成クラス"""
    
    def __init__(self, output_folder):
        self.output_folder = output_folder
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
    def create_cover_slide(self, summary_stats, video_duration):
        """カバースライド作成"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # 背景色設定
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # タイトル
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = "ライブコマース配信分析レポート"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(51, 51, 51)
        title_para.alignment = PP_ALIGN.CENTER
        
        # サブタイトル
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4),
            Inches(8), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"対象データ: 1〜{video_duration}分 / 生成日: {datetime.now().strftime('%Y年%m月%d日')}"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.color.rgb = RGBColor(102, 126, 234)
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def create_kpi_summary_slide(self, summary_stats):
        """主要KPIサマリースライド作成"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # タイトル
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "主要KPIサマリー"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(51, 51, 51)
        
        # KPIカード配置
        kpis = [
            {
                'title': '最大同時視聴者数',
                'value': f"{summary_stats.get('max_viewers', 0):,}",
                'unit': '人',
                'color': RGBColor(102, 126, 234)
            },
            {
                'title': '合計いいね数',
                'value': f"{summary_stats.get('total_likes', 0):,}",
                'unit': '件',
                'color': RGBColor(233, 30, 99)
            },
            {
                'title': '合計コメント数',
                'value': f"{summary_stats.get('total_comments_actual', summary_stats.get('total_comments_metric', 0)):,}",
                'unit': '件',
                'color': RGBColor(76, 175, 80)
            },
            {
                'title': '合計クリック数',
                'value': f"{summary_stats.get('total_clicks', 0):,}",
                'unit': '件',
                'color': RGBColor(255, 152, 0)
            }
        ]
        
        # 2行x2列のグリッド
        row = 0
        col = 0
        for kpi in kpis:
            left = Inches(0.5 + col * 4.5)
            top = Inches(1.5 + row * 2.5)
            width = Inches(4)
            height = Inches(2)
            
            # KPIカード作成
            self._create_kpi_card(slide, left, top, width, height, kpi)
            
            col += 1
            if col >= 2:
                col = 0
                row += 1
        
        return slide
    
    def _create_kpi_card(self, slide, left, top, width, height, kpi):
        """KPIカード作成"""
        # 背景
        card = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, width, height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(248, 249, 250)
        card.line.color.rgb = kpi['color']
        card.line.width = Pt(2)
        
        # タイトル
        title_box = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.2),
            width - Inches(0.4), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = kpi['title']
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(14)
        title_para.font.color.rgb = RGBColor(102, 102, 102)
        
        # 値
        value_box = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.8),
            width - Inches(0.4), Inches(0.8)
        )
        value_frame = value_box.text_frame
        value_frame.text = f"{kpi['value']} {kpi['unit']}"
        value_para = value_frame.paragraphs[0]
        value_para.font.size = Pt(36)
        value_para.font.bold = True
        value_para.font.color.rgb = kpi['color']
    
    def create_chart_slide(self, title, chart_image_path):
        """グラフスライド作成"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # タイトル
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(51, 51, 51)
        
        # グラフ画像
        if os.path.exists(chart_image_path):
            slide.shapes.add_picture(
                chart_image_path,
                Inches(0.5), Inches(1.5),
                width=Inches(9), height=Inches(5.5)
            )
        
        return slide
    
    def create_comment_analysis_slide(self, comment_analysis, pie_chart_path):
        """コメント分析スライド作成"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # タイトル
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "コメント分類分析"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(51, 51, 51)
        
        # 円グラフ
        if os.path.exists(pie_chart_path):
            slide.shapes.add_picture(
                pie_chart_path,
                Inches(0.5), Inches(1.5),
                width=Inches(5), height=Inches(5.5)
            )
        
        # カテゴリ詳細
        categories = comment_analysis.get('categories', {})
        left = Inches(5.8)
        top = Inches(1.5)
        
        for i, (category, count) in enumerate(categories.items()):
            # カテゴリ名
            cat_box = slide.shapes.add_textbox(
                left, top + Inches(i * 0.8),
                Inches(3.5), Inches(0.35)
            )
            cat_frame = cat_box.text_frame
            cat_frame.text = f"{category}: {count}件"
            cat_para = cat_frame.paragraphs[0]
            cat_para.font.size = Pt(16)
            cat_para.font.bold = True
            cat_para.font.color.rgb = RGBColor(51, 51, 51)
        
        return slide
    
    def create_recommendations_slide(self, recommendations):
        """改善提案スライド作成"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # タイトル
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "改善提案"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(51, 51, 51)
        
        # Good Points
        self._add_recommendation_section(
            slide,
            Inches(0.5), Inches(1.5),
            "✅ 良かった点 (Good)",
            recommendations.get('good_points', []),
            RGBColor(212, 237, 218)
        )
        
        # Improvements
        self._add_recommendation_section(
            slide,
            Inches(0.5), Inches(3.5),
            "📈 改善すべき点 (More)",
            recommendations.get('improvements', []),
            RGBColor(255, 243, 205)
        )
        
        # Next Actions
        self._add_recommendation_section(
            slide,
            Inches(0.5), Inches(5.5),
            "🎬 次回に向けたアクション",
            recommendations.get('next_actions', [])[:2],  # 2件まで
            RGBColor(209, 236, 241)
        )
        
        return slide
    
    def _add_recommendation_section(self, slide, left, top, title, items, bg_color):
        """推奨事項セクション追加"""
        width = Inches(9)
        height = Inches(1.8)
        
        # 背景
        bg = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, width, height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.width = Pt(0)
        
        # タイトル
        title_box = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.1),
            width - Inches(0.4), Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(51, 51, 51)
        
        # 項目
        items_box = slide.shapes.add_textbox(
            left + Inches(0.3), top + Inches(0.6),
            width - Inches(0.6), Inches(1)
        )
        items_frame = items_box.text_frame
        items_frame.word_wrap = True
        
        for item in items[:2]:  # 最大2件
            p = items_frame.add_paragraph()
            p.text = f"• {item[:80]}..."  # 80文字まで
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.space_after = Pt(6)
    
    def save(self, filename="report.pptx"):
        """PPTXファイルを保存"""
        filepath = os.path.join(self.output_folder, filename)
        self.prs.save(filepath)
        return filepath
