"""
強化版PowerPoint生成クラス
サンプルレポートの構成・見せ方を学習して実装
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
from datetime import datetime


class EnhancedPowerPointGenerator:
    """サンプルレポートに基づく強化版PowerPoint生成クラス"""
    
    def __init__(self, output_folder):
        """初期化"""
        self.output_folder = output_folder
        self.prs = Presentation()
        
        # スライドサイズ（16:9ワイドスクリーン）
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        
        # カラーパレット（サンプルレポートに合わせた配色）
        self.colors = {
            'primary_blue': RGBColor(68, 114, 196),      # メインブルー
            'accent_blue': RGBColor(91, 155, 213),       # アクセントブルー
            'light_blue': RGBColor(221, 235, 247),       # 薄いブルー（背景用）
            'dark_text': RGBColor(51, 51, 51),           # 濃いグレー（本文）
            'light_text': RGBColor(127, 127, 127),       # 薄いグレー（補足）
            'white': RGBColor(255, 255, 255),            # 白
            'success_green': RGBColor(146, 208, 80),     # 成功指標用
            'warning_orange': RGBColor(255, 192, 0),     # 注意喚起用
            'metric_red': RGBColor(255, 102, 102),       # 重要指標用
        }
    
    def create_slide_1_cover(self, summary_stats, video_duration):
        """スライド1: カバーページ"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 背景（薄いブルー）
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['light_blue']
        
        # メインタイトル
        title = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(11.33), Inches(1)
        )
        tf = title.text_frame
        tf.text = "ライブコマース配信分析レポート"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        
        # サブタイトル
        subtitle = slide.shapes.add_textbox(
            Inches(1), Inches(3.7),
            Inches(11.33), Inches(0.6)
        )
        tf = subtitle.text_frame
        tf.text = "データドリブン分析による売上最大化インサイト"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.color.rgb = self.colors['dark_text']
        
        # 配信概要
        info = slide.shapes.add_textbox(
            Inches(1), Inches(4.8),
            Inches(11.33), Inches(0.5)
        )
        tf = info.text_frame
        tf.text = f"対象データ: 開始1〜{video_duration}分 ／ 最大視聴者数: {summary_stats.get('max_viewers', 0):,}人"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.color.rgb = self.colors['light_text']
        
        # 生成日時
        date_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5),
            Inches(11.33), Inches(0.4)
        )
        tf = date_box.text_frame
        tf.text = f"レポート生成日: {datetime.now().strftime('%Y年%m月%d日')}"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.color.rgb = self.colors['light_text']
    
    def create_slide_2_kpi_summary(self, summary_stats, peak_analysis):
        """スライド2: 主要KPIサマリー"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # タイトル
        self._add_slide_title(slide, "主要KPIサマリー")
        
        # サブタイトル
        subtitle = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.95),
            Inches(12.33), Inches(0.3)
        )
        tf = subtitle.text_frame
        tf.text = "配信全体のパフォーマンス指標と主要ピーク"
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = self.colors['light_text']
        
        # KPIカード（2行x3列レイアウト）
        kpis = [
            {
                'title': '最大同時視聴者数',
                'value': f"{summary_stats.get('max_viewers', 0):,}",
                'unit': '人',
                'peak_info': self._get_peak_minute(peak_analysis, 'viewers'),
                'color': self.colors['primary_blue']
            },
            {
                'title': '平均視聴者数',
                'value': f"{summary_stats.get('avg_viewers', 0):.0f}",
                'unit': '人',
                'peak_info': '',
                'color': self.colors['accent_blue']
            },
            {
                'title': '合計いいね数',
                'value': f"{summary_stats.get('total_likes', 0):,}",
                'unit': '件',
                'peak_info': self._get_peak_minute(peak_analysis, 'likes'),
                'color': self.colors['metric_red']
            },
            {
                'title': '合計コメント数',
                'value': f"{summary_stats.get('total_comments_actual', summary_stats.get('total_comments_metric', 0)):,}",
                'unit': '件',
                'peak_info': self._get_peak_minute(peak_analysis, 'comments'),
                'color': self.colors['success_green']
            },
            {
                'title': '合計クリック数',
                'value': f"{summary_stats.get('total_clicks', 0):,}",
                'unit': '件',
                'peak_info': self._get_peak_minute(peak_analysis, 'clicks'),
                'color': self.colors['warning_orange']
            },
            {
                'title': 'エンゲージメント率',
                'value': self._calculate_engagement_rate(summary_stats),
                'unit': '%',
                'peak_info': '(いいね+コメント)/視聴者',
                'color': self.colors['primary_blue']
            }
        ]
        
        # KPIカード配置
        start_y = 1.5
        card_width = 3.8
        card_height = 2.2
        x_gap = 0.3
        y_gap = 0.3
        
        for idx, kpi in enumerate(kpis):
            row = idx // 3
            col = idx % 3
            
            x = 0.5 + col * (card_width + x_gap)
            y = start_y + row * (card_height + y_gap)
            
            self._create_kpi_card_enhanced(slide, x, y, card_width, card_height, kpi)
    
    def create_slide_3_timeline_viewers(self, chart_path, peak_info):
        """スライド3: 時系列(1) 同時視聴ユーザー数の推移"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # タイトル
        self._add_slide_title(slide, "時系列(1) 同時視聴ユーザー数の推移")
        
        # グラフ画像
        if os.path.exists(chart_path):
            slide.shapes.add_picture(
                chart_path,
                Inches(0.5), Inches(1.5),
                width=Inches(8.5), height=Inches(5.2)
            )
        
        # 注釈ボックス
        annotation = slide.shapes.add_textbox(
            Inches(9.2), Inches(1.5),
            Inches(3.8), Inches(5.2)
        )
        tf = annotation.text_frame
        tf.word_wrap = True
        
        # 注釈タイトル
        p = tf.paragraphs[0]
        p.text = "📊 主要インサイト"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        p.space_after = Pt(10)
        
        # ピーク情報
        if peak_info and 'viewers' in peak_info and peak_info['viewers']:
            for peak in peak_info['viewers'][:3]:
                p = tf.add_paragraph()
                p.text = f"• {peak['minute']}分: {peak['value']:.0f}人"
                p.font.size = Pt(12)
                p.space_after = Pt(8)
                
                # 説明
                p = tf.add_paragraph()
                p.text = f"  {peak.get('event_description', '増加トレンド確認')}"
                p.font.size = Pt(10)
                p.font.color.rgb = self.colors['light_text']
                p.space_after = Pt(12)
    
    def create_slide_4_timeline_clicks(self, chart_path, peak_info):
        """スライド4: 時系列(2) 商品クリック数とカート追加"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_slide_title(slide, "時系列(2) 商品クリック数の推移")
        
        # グラフ
        if os.path.exists(chart_path):
            slide.shapes.add_picture(
                chart_path,
                Inches(0.5), Inches(1.5),
                width=Inches(8.5), height=Inches(5.2)
            )
        
        # 注釈
        annotation = slide.shapes.add_textbox(
            Inches(9.2), Inches(1.5),
            Inches(3.8), Inches(5.2)
        )
        tf = annotation.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "🔍 クリック動向分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        p.space_after = Pt(10)
        
        if peak_info and 'clicks' in peak_info and peak_info['clicks']:
            for peak in peak_info['clicks'][:3]:
                p = tf.add_paragraph()
                p.text = f"• {peak['minute']}分: {peak['value']:.0f}件"
                p.font.size = Pt(12)
                p.space_after = Pt(8)
                
                p = tf.add_paragraph()
                p.text = f"  {peak.get('event_description', '商品紹介効果')}"
                p.font.size = Pt(10)
                p.font.color.rgb = self.colors['light_text']
                p.space_after = Pt(12)
    
    def create_slide_5_timeline_engagement(self, chart_path, peak_info):
        """スライド5: 時系列(3) いいね数とチャット数"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_slide_title(slide, "時系列(3) いいね数とコメント数")
        
        if os.path.exists(chart_path):
            slide.shapes.add_picture(
                chart_path,
                Inches(0.5), Inches(1.5),
                width=Inches(8.5), height=Inches(5.2)
            )
        
        # 注釈
        annotation = slide.shapes.add_textbox(
            Inches(9.2), Inches(1.5),
            Inches(3.8), Inches(5.2)
        )
        tf = annotation.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "💬 エンゲージメント分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        p.space_after = Pt(10)
        
        # いいねピーク
        if peak_info and 'likes' in peak_info and peak_info['likes']:
            p = tf.add_paragraph()
            p.text = "【いいね】"
            p.font.size = Pt(12)
            p.font.bold = True
            p.space_after = Pt(6)
            
            for peak in peak_info['likes'][:2]:
                p = tf.add_paragraph()
                p.text = f"• {peak['minute']}分: {peak['value']:.0f}件"
                p.font.size = Pt(11)
                p.space_after = Pt(6)
        
        # コメントピーク
        if peak_info and 'comments' in peak_info and peak_info['comments']:
            p = tf.add_paragraph()
            p.text = "\n【コメント】"
            p.font.size = Pt(12)
            p.font.bold = True
            p.space_after = Pt(6)
            
            for peak in peak_info['comments'][:2]:
                p = tf.add_paragraph()
                p.text = f"• {peak['minute']}分: {peak['value']:.0f}件"
                p.font.size = Pt(11)
                p.space_after = Pt(6)
    
    def create_slide_6_single_metric_viewers(self, peak_info, recommendations):
        """スライド6: 単一指標分析｜同時視聴ユーザー数"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_slide_title(slide, "単一指標分析｜同時視聴ユーザー数")
        
        # 3カラムレイアウト
        self._create_three_column_analysis(
            slide,
            title="視聴者維持の鍵",
            peak_label="ピーク",
            peak_value=self._get_max_peak_value(peak_info, 'viewers'),
            insights_title="洞察",
            insights=self._extract_viewer_insights(peak_info, recommendations),
            improvements_title="改善施策",
            improvements=self._extract_viewer_improvements(recommendations)
        )
    
    def create_slide_7_single_metric_clicks(self, peak_info, recommendations):
        """スライド7: 単一指標分析｜商品クリック数"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_slide_title(slide, "単一指標分析｜商品クリック数")
        
        self._create_three_column_analysis(
            slide,
            title="購買行動の促進",
            peak_label="ピーク",
            peak_value=self._get_max_peak_value(peak_info, 'clicks'),
            insights_title="洞察",
            insights=self._extract_click_insights(peak_info, recommendations),
            improvements_title="改善施策",
            improvements=self._extract_click_improvements(recommendations)
        )
    
    def create_slide_8_single_metric_engagement(self, peak_info, recommendations):
        """スライド8: 単一指標分析｜チャット＆いいね"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_slide_title(slide, "単一指標分析｜チャット＆いいね")
        
        likes_peak = self._get_max_peak_value(peak_info, 'likes')
        comments_peak = self._get_max_peak_value(peak_info, 'comments')
        
        self._create_three_column_analysis(
            slide,
            title="双方向コミュニケーション",
            peak_label="ピーク",
            peak_value=f"いいね{likes_peak} / コメント{comments_peak}",
            insights_title="洞察",
            insights=self._extract_engagement_insights(peak_info, recommendations),
            improvements_title="改善施策",
            improvements=self._extract_engagement_improvements(recommendations)
        )
    
    def create_slide_9_multi_metric_correlation(self, summary_stats, peak_info, recommendations):
        """スライド9: 複数指標分析｜視聴×クリックの相関"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_slide_title(slide, "複数指標分析｜視聴×クリックの相関と課題")
        
        # CTR計算
        ctr = self._calculate_ctr(summary_stats)
        
        self._create_three_column_analysis(
            slide,
            title="コンバージョン分析",
            peak_label="推定CTR",
            peak_value=f"{ctr:.2f}%",
            insights_title="相関と示唆",
            insights=self._extract_correlation_insights(summary_stats, peak_info),
            improvements_title="対策",
            improvements=self._extract_correlation_improvements(recommendations)
        )
    
    def create_slide_10_comment_analysis(self, comment_analysis, pie_chart_path):
        """スライド10: コメント定量分析"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_slide_title(slide, "コメント定量分析")
        
        # 左側: 円グラフ
        if os.path.exists(pie_chart_path):
            slide.shapes.add_picture(
                pie_chart_path,
                Inches(0.8), Inches(1.8),
                width=Inches(5.5), height=Inches(4.8)
            )
        
        # 右側: カテゴリ詳細
        detail_box = slide.shapes.add_textbox(
            Inches(6.8), Inches(1.5),
            Inches(6), Inches(5.5)
        )
        tf = detail_box.text_frame
        tf.word_wrap = True
        
        # タイトル
        p = tf.paragraphs[0]
        p.text = "カテゴリ別内訳と示唆"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        p.space_after = Pt(12)
        
        # カテゴリ詳細
        categories = comment_analysis.get('categories', {})
        total = comment_analysis.get('total', 0)
        
        for category, count in sorted(categories.items(), key=lambda x: x[1], reverse=True):
            percentage = (count / total * 100) if total > 0 else 0
            
            p = tf.add_paragraph()
            p.text = f"• {category}: {count}件 ({percentage:.1f}%)"
            p.font.size = Pt(13)
            p.font.bold = True
            p.space_after = Pt(4)
            
            # カテゴリごとの示唆
            insight = self._get_category_insight(category, percentage)
            if insight:
                p = tf.add_paragraph()
                p.text = f"  → {insight}"
                p.font.size = Pt(11)
                p.font.color.rgb = self.colors['light_text']
                p.space_after = Pt(10)
    
    def create_slide_11_overall_insights(self, recommendations, summary_stats):
        """スライド11: 総合考察｜成功要因と課題"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_slide_title(slide, "総合考察｜成功要因と課題")
        
        # 2カラムレイアウト
        # 左カラム: 成功要因
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(6), Inches(5.5)
        )
        tf = left_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "✅ 成功要因"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['success_green']
        p.space_after = Pt(12)
        
        good_points = recommendations.get('good_points', [])
        for idx, point in enumerate(good_points[:5], 1):
            p = tf.add_paragraph()
            p.text = f"{idx}. {point}"
            p.font.size = Pt(12)
            p.space_after = Pt(10)
            p.level = 0
        
        # 右カラム: 課題
        right_box = slide.shapes.add_textbox(
            Inches(7), Inches(1.5),
            Inches(6), Inches(5.5)
        )
        tf = right_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "⚠️ 改善すべき課題"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['warning_orange']
        p.space_after = Pt(12)
        
        improvements = recommendations.get('improvements', [])
        for idx, improvement in enumerate(improvements[:5], 1):
            p = tf.add_paragraph()
            p.text = f"{idx}. {improvement}"
            p.font.size = Pt(12)
            p.space_after = Pt(10)
            p.level = 0
    
    def create_slide_12_action_plan(self, recommendations):
        """スライド12: アクションプラン（次回配信）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_slide_title(slide, "アクションプラン（次回配信）")
        
        # 次回アクション
        action_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5),
            Inches(11.73), Inches(5.5)
        )
        tf = action_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "🎯 次回配信に向けた改善施策"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        p.space_after = Pt(16)
        
        next_actions = recommendations.get('next_actions', [])
        for idx, action in enumerate(next_actions, 1):
            # アクション番号
            p = tf.add_paragraph()
            p.text = f"{idx}"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.space_before = Pt(12)
            p.space_after = Pt(8)
            
            # 背景色付きの番号（疑似的に表現）
            # （実際の背景色はテキストボックスでは制限があるため、シンプルに表現）
            
            # アクション内容
            p = tf.add_paragraph()
            p.text = action
            p.font.size = Pt(14)
            p.space_after = Pt(16)
            p.level = 0
        
        # フッター
        footer = slide.shapes.add_textbox(
            Inches(0.8), Inches(6.8),
            Inches(11.73), Inches(0.4)
        )
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.text = "💡 ヒント: これらの施策を1つずつ実践し、次回配信でA/Bテストを実施することをお勧めします"
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = self.colors['light_text']
    
    # ヘルパーメソッド
    
    def _add_slide_title(self, slide, title_text):
        """スライドタイトルを追加"""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.33), Inches(0.6)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
    
    def _create_kpi_card_enhanced(self, slide, left, top, width, height, kpi):
        """強化版KPIカード作成"""
        # 背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.colors['white']
        card.line.color.rgb = kpi['color']
        card.line.width = Pt(3)
        
        # タイトル
        title_box = slide.shapes.add_textbox(
            Inches(left + 0.2), Inches(top + 0.2),
            Inches(width - 0.4), Inches(0.4)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = kpi['title']
        p.font.size = Pt(13)
        p.font.color.rgb = self.colors['dark_text']
        
        # 値
        value_box = slide.shapes.add_textbox(
            Inches(left + 0.2), Inches(top + 0.7),
            Inches(width - 0.4), Inches(0.8)
        )
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{kpi['value']} {kpi['unit']}"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = kpi['color']
        
        # ピーク情報
        if kpi.get('peak_info'):
            peak_box = slide.shapes.add_textbox(
                Inches(left + 0.2), Inches(top + 1.6),
                Inches(width - 0.4), Inches(0.4)
            )
            tf = peak_box.text_frame
            p = tf.paragraphs[0]
            p.text = kpi['peak_info']
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors['light_text']
    
    def _create_three_column_analysis(self, slide, title, peak_label, peak_value, 
                                      insights_title, insights, improvements_title, improvements):
        """3カラム分析レイアウト作成"""
        # サブタイトル
        subtitle = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.95),
            Inches(12.33), Inches(0.3)
        )
        tf = subtitle.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.color.rgb = self.colors['light_text']
        
        col_width = 3.9
        col_height = 5.2
        start_y = 1.5
        gap = 0.3
        
        # 左カラム: ピーク
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(start_y),
            Inches(col_width), Inches(col_height)
        )
        tf = left_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = peak_label
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        p.space_after = Pt(12)
        
        p = tf.add_paragraph()
        p.text = peak_value
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        
        # 中央カラム: 洞察
        center_box = slide.shapes.add_textbox(
            Inches(0.5 + col_width + gap), Inches(start_y),
            Inches(col_width), Inches(col_height)
        )
        tf = center_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = insights_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.colors['success_green']
        p.space_after = Pt(12)
        
        for idx, insight in enumerate(insights[:4], 1):
            p = tf.add_paragraph()
            p.text = f"{idx}. {insight}"
            p.font.size = Pt(11)
            p.space_after = Pt(10)
        
        # 右カラム: 改善施策
        right_box = slide.shapes.add_textbox(
            Inches(0.5 + 2 * (col_width + gap)), Inches(start_y),
            Inches(col_width), Inches(col_height)
        )
        tf = right_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = improvements_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.colors['warning_orange']
        p.space_after = Pt(12)
        
        for idx, improvement in enumerate(improvements[:4], 1):
            p = tf.add_paragraph()
            p.text = f"{idx}. {improvement}"
            p.font.size = Pt(11)
            p.space_after = Pt(10)
    
    def _get_peak_minute(self, peak_analysis, metric):
        """ピーク時刻を取得"""
        if metric in peak_analysis and peak_analysis[metric]:
            peak = peak_analysis[metric][0]
            return f"{peak['minute']}分でピーク"
        return ""
    
    def _get_max_peak_value(self, peak_analysis, metric):
        """最大ピーク値を取得"""
        if metric in peak_analysis and peak_analysis[metric]:
            peak = peak_analysis[metric][0]
            return f"{peak['value']:.0f}({peak['minute']}分)"
        return "N/A"
    
    def _calculate_engagement_rate(self, summary_stats):
        """エンゲージメント率を計算"""
        total_engagement = summary_stats.get('total_likes', 0) + summary_stats.get('total_comments_actual', 0)
        max_viewers = summary_stats.get('max_viewers', 1)
        if max_viewers > 0:
            return f"{(total_engagement / max_viewers * 100):.1f}"
        return "0.0"
    
    def _calculate_ctr(self, summary_stats):
        """CTRを計算"""
        clicks = summary_stats.get('total_clicks', 0)
        viewers = summary_stats.get('max_viewers', 1)
        if viewers > 0:
            return (clicks / viewers) * 100
        return 0.0
    
    def _extract_viewer_insights(self, peak_info, recommendations):
        """視聴者数の洞察を抽出"""
        insights = []
        if 'viewers' in peak_info and peak_info['viewers']:
            for peak in peak_info['viewers'][:2]:
                insights.append(f"{peak['minute']}分に{peak['value']:.0f}人のピークを記録")
        
        # 推奨事項から関連する洞察を追加
        good_points = recommendations.get('good_points', [])
        for point in good_points[:2]:
            if '視聴' in point or 'ユーザー' in point:
                insights.append(point[:60])
        
        if not insights:
            insights.append("視聴者の関心を維持するための施策が必要")
        
        return insights
    
    def _extract_viewer_improvements(self, recommendations):
        """視聴者数の改善案を抽出"""
        improvements = recommendations.get('improvements', [])
        viewer_improvements = [imp for imp in improvements if '視聴' in imp or '離脱' in imp]
        
        if not viewer_improvements:
            viewer_improvements = [
                "冒頭30秒で視聴者の関心を引く工夫",
                "定期的な視聴者への呼びかけ",
                "視聴維持のための起伏ある構成"
            ]
        
        return viewer_improvements[:4]
    
    def _extract_click_insights(self, peak_info, recommendations):
        """クリック数の洞察を抽出"""
        insights = []
        if 'clicks' in peak_info and peak_info['clicks']:
            for peak in peak_info['clicks'][:2]:
                insights.append(f"{peak['minute']}分に{peak['value']:.0f}件のクリック")
        
        good_points = recommendations.get('good_points', [])
        for point in good_points:
            if 'クリック' in point or '商品' in point:
                insights.append(point[:60])
                break
        
        if not insights:
            insights.append("商品紹介のタイミングと方法の最適化が必要")
        
        return insights
    
    def _extract_click_improvements(self, recommendations):
        """クリック数の改善案を抽出"""
        improvements = recommendations.get('improvements', [])
        click_improvements = [imp for imp in improvements if 'クリック' in imp or '商品' in imp]
        
        if not click_improvements:
            click_improvements = [
                "商品の魅力を視覚的に訴求",
                "限定性・緊急性の演出",
                "具体的な使用シーンの提示"
            ]
        
        return click_improvements[:4]
    
    def _extract_engagement_insights(self, peak_info, recommendations):
        """エンゲージメントの洞察を抽出"""
        insights = []
        
        if 'likes' in peak_info and peak_info['likes']:
            peak = peak_info['likes'][0]
            insights.append(f"いいね: {peak['minute']}分に{peak['value']:.0f}件")
        
        if 'comments' in peak_info and peak_info['comments']:
            peak = peak_info['comments'][0]
            insights.append(f"コメント: {peak['minute']}分に{peak['value']:.0f}件")
        
        good_points = recommendations.get('good_points', [])
        for point in good_points:
            if 'コメント' in point or 'いいね' in point or 'エンゲージ' in point:
                insights.append(point[:60])
                break
        
        return insights or ["双方向コミュニケーションの強化が必要"]
    
    def _extract_engagement_improvements(self, recommendations):
        """エンゲージメントの改善案を抽出"""
        improvements = recommendations.get('improvements', [])
        eng_improvements = [imp for imp in improvements if 'コメント' in imp or 'いいね' in imp]
        
        if not eng_improvements:
            eng_improvements = [
                "視聴者への質問投げかけ",
                "コメントへのリアルタイム反応",
                "参加型企画の実施"
            ]
        
        return eng_improvements[:4]
    
    def _extract_correlation_insights(self, summary_stats, peak_info):
        """相関分析の洞察を抽出"""
        ctr = self._calculate_ctr(summary_stats)
        insights = [
            f"推定CTR: {ctr:.2f}%",
            "視聴者数とクリック数に一定の相関",
        ]
        
        if ctr > 30:
            insights.append("高いCTRは商品訴求力の証")
        elif ctr > 15:
            insights.append("平均的なCTR、改善の余地あり")
        else:
            insights.append("CTR向上のための施策が急務")
        
        return insights
    
    def _extract_correlation_improvements(self, recommendations):
        """相関分析の改善案を抽出"""
        improvements = recommendations.get('next_actions', [])
        
        if not improvements:
            improvements = [
                "商品紹介とCTAのタイミング最適化",
                "視覚的訴求力の強化",
                "購入導線の簡略化"
            ]
        
        return improvements[:4]
    
    def _get_category_insight(self, category, percentage):
        """カテゴリ別の示唆を生成"""
        insights = {
            '質問': "視聴者の疑問に即座に回答することで信頼構築" if percentage > 20 else "Q&Aタイムの設定を検討",
            '驚き': "驚きの要素が視聴継続を促進" if percentage > 10 else "サプライズ要素の追加を検討",
            'ワクワク・期待感': "期待感の醸成が成功" if percentage > 15 else "ストーリーテリングの強化",
            '挨拶': "コミュニティ感の醸成" if percentage < 30 else "挨拶が多すぎる可能性",
            '購入意志': "高い購買意欲を確認" if percentage > 15 else "購買促進施策が必要",
            'その他': "多様なリアクション" if percentage < 40 else "明確な反応を引き出す工夫"
        }
        return insights.get(category, "")
    
    def save(self, filename="report.pptx"):
        """PPTXファイルを保存"""
        filepath = os.path.join(self.output_folder, filename)
        self.prs.save(filepath)
        return filepath
